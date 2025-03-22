import pdfplumber
from pptx import Presentation
from docx import Document


def extract_text(file):
    if file.content_type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file)
    elif file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_docx(file)
    
    return ""
        

def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text


def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text


def extract_text_from_docx(file):
    doc = Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text