import pdfplumber
from pptx import Presentation
from docx import Document
from google import generativeai as genai
from dotenv import load_dotenv
import os

load_dotenv()
genai.configure(api_key=os.getenv("GENAI_API_KEY"))


def extract_text(file):
    if file.content_type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file)
    elif file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_docx(file)
    
    return ""
        

def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text()
    return text


def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text


def extract_text_from_docx(file):
    doc = Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text


def extract_title(file):
    text = extract_text(file)
    model = genai.GenerativeModel(model_name='gemini-2.0-flash')
    response = model.generate_content(f"Extract the lecture title and respond with only the title and nothing else. Do not include course code and make sure the title is coherent. Do not append newline to your response:\n\nDocument:\n{text}")
    return response.text
